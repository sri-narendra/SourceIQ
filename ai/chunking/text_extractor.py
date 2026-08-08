import io

# Googsheet-style plain-text formats we can safely decode as text.
TEXT_EXTENSIONS = {
    "txt",
    "md",
    "markdown",
    "csv",
    "tsv",
    "json",
    "xml",
    "html",
    "htm",
    "yaml",
    "yml",
    "toml",
    "ini",
    "cfg",
    "log",
    "rtf",
    "py",
    "js",
    "ts",
    "jsx",
    "tsx",
    "css",
    "scss",
    "sql",
    "sh",
    "bash",
    "java",
    "c",
    "cpp",
    "h",
    "hpp",
    "go",
    "rs",
    "rb",
    "php",
    "swift",
    "kt",
    "scala",
    "r",
}

# Formats with dedicated extractors.
BINARY_EXTENSIONS = {"pdf", "docx", "pptx", "xlsx", "odt"}

# Raster images: always OCR'd.
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "bmp", "tif", "tiff", "webp"}

SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | BINARY_EXTENSIONS | IMAGE_EXTENSIONS


class UnknownTypeError(Exception):
    pass


def extract_text(filename: str, data: bytes) -> str:
    pages = extract_pages(filename, data)
    return "\n\n".join(text for text, _ in pages)


def extract_pages(filename: str, data: bytes) -> list[tuple[str, int | None]]:
    """Extract text per page where the format has pages; else one page.

    Returns [(text, page_number)] with 1-based page numbers for PDFs, and
    page=None for formats with no page concept (plain text, office docs).
    Scanned PDFs fall back to per-page OCR.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in TEXT_EXTENSIONS:
        return [(data.decode("utf-8", errors="ignore"), None)]
    if ext == "pdf":
        return _pdf_pages(data)
    if ext in ("docx", "pptx", "xlsx", "odt"):
        return [(_office(data, ext), None)]
    if ext in IMAGE_EXTENSIONS:
        return [(_ocr_bytes(data), None)]
    raise UnknownTypeError(f"Unsupported file type: {ext}")


def _pdf_pages(data: bytes) -> list[tuple[str, int]]:
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [(page.extract_text() or "", i + 1) for i, page in enumerate(reader.pages)]
    if sum(len(t) for t, _ in pages) < 20:
        # ponytail: scanned PDFs have no text layer — render pages and OCR;
        # threshold 20 chars catches the common "blank" case without extra cost.
        return _pdf_ocr_pages(data)
    return [(t.strip(), p) for t, p in pages]


def _pdf_ocr_pages(data: bytes) -> list[tuple[str, int]]:
    import pymupdf

    doc = pymupdf.open(stream=data)
    return [
        (_ocr_bytes(page.get_pixmap(dpi=200).tobytes("png")), i + 1)
        for i, page in enumerate(doc)
    ]


def _ocr_bytes(data: bytes) -> str:
    from rapidocr_onnxruntime import RapidOCR

    # ponytail: one shared engine; cached per-process by rapidocr.
    result, _ = RapidOCR()(data)
    if not result:
        return ""
    return "\n".join(line[1] for line in result)


def _office(data: bytes, ext: str) -> str:
    if ext == "docx":
        import docx

        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    if ext == "pptx":
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        lines.append("".join(run.text for run in para.runs))
        return "\n".join(lines)
    if ext == "xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"[{sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                lines.append(", ".join("" if c is None else str(c) for c in row))
        return "\n".join(lines)
    if ext == "odt":
        import zipfile
        from xml.etree import ElementTree

        with zipfile.ZipFile(io.BytesIO(data)) as z:
            xml = z.read("content.xml")
        ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
        try:
            root = ElementTree.fromstring(xml)
        except Exception:
            return ""
        return "\n".join(
            (p.text or "") for p in root.iter("{urn:oasis:names:tc:opendocument:xmlns:text:1.0}p")
        )
    raise UnknownTypeError(f"Unsupported file type: {ext}")